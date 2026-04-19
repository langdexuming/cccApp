from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


STATUS_LABELS = {
    "STARTED": "进行中",
    "COMPLETED": "已完成",
    "SKIPPED": "已跳过",
    "PLANNED": "预留",
    "BLOCKED": "阻塞",
    "FAILED": "失败",
}

STATUS_COLORS = {
    "STARTED": "#4CC9F0",
    "COMPLETED": "#80ED99",
    "SKIPPED": "#A0AEC0",
    "PLANNED": "#F6C453",
    "BLOCKED": "#F59E0B",
    "FAILED": "#F87171",
}

STEP_LABELS = {
    "RUN": "运行",
    "PREPARE": "准备",
    "TRAIN": "训练",
    "EXPORT": "导出",
    "GGUF": "GGUF 转换",
    "OLLAMA_REGISTER": "Ollama 注册",
    "VLLM_REGISTER": "vLLM 注册",
    "EVAL": "评测",
    "BENCHMARK": "基准压测",
    "REPORT": "报告",
}


@dataclass(slots=True)
class PipelineEntry:
    timestamp: str
    step: str
    status: str
    detail: str


def find_training_intro_ppt(resources_dir: Path) -> Path:
    exact = resources_dir / "Qwen3.5-9B-Base动环告警诊断训练介绍.pptx"
    if exact.exists():
        return exact
    return next(
        p
        for p in resources_dir.glob("*.pptx")
        if "Qwen3.5-9B-Base" in p.name
        and "bak_fillshots" not in p.name
        and "copytest" not in p.name
        and not p.name.startswith("~$")
    )


def safe_read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def load_json(path: Path) -> dict[str, Any]:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def parse_exit_code(log_path: Path) -> int | None:
    if not log_path.exists():
        return None
    match = re.search(r"\[exit_code\]\s+(-?\d+)", safe_read_text(log_path))
    return int(match.group(1)) if match else None


def fmt_time(dt: datetime) -> str:
    return dt.strftime("%Y-%m-%d %H:%M:%S")


def next_time(current: datetime | None, fallback: datetime) -> datetime:
    if current is None:
        return fallback
    return current + timedelta(seconds=16)


def coerce_dt(path: Path, default: datetime) -> datetime:
    if path.exists():
        detected = datetime.fromtimestamp(path.stat().st_mtime)
        return detected if detected >= default else default
    return default


def build_entries_for_run(run_dir: Path) -> list[PipelineEntry]:
    now = datetime.now()
    input_dir = run_dir / "input"
    output_dir = run_dir / "output"
    log_dir = run_dir / "logs"
    eval_dir = run_dir / "eval"
    benchmark_dir = run_dir / "benchmark"
    report_path = run_dir.parent.parent / "reports" / "single" / f"{run_dir.name}.json"

    manifest = load_json(input_dir / "run_manifest.json")
    preview = load_json(output_dir / "command_preview.json")
    eval_summary = load_json(eval_dir / "eval_summary.json")
    benchmark_summary = load_json(benchmark_dir / "benchmark_summary.json")

    route_type = str(manifest.get("route_type") or "unknown")
    trainer_backend = str(manifest.get("trainer_backend") or "unknown")
    base_model = str(manifest.get("base_model") or "unknown")
    preview_meta = preview.get("meta") or {}

    entries: list[PipelineEntry] = []
    cursor = coerce_dt(input_dir / "run_manifest.json", now)
    entries.append(
        PipelineEntry(
            timestamp=fmt_time(cursor),
            step="RUN",
            status="STARTED",
            detail=f"run_id={run_dir.name} trainer_backend={trainer_backend} base_model={base_model}",
        )
    )

    cursor = next_time(cursor, now)
    if preview:
        entries.append(
            PipelineEntry(
                timestamp=fmt_time(cursor),
                step="PREPARE",
                status="COMPLETED",
                detail=f"command_preview={output_dir / 'command_preview.json'}",
            )
        )

    train_log = log_dir / "train.log"
    adapter_dir = output_dir / "adapter"
    cursor = next_time(cursor, now)
    if train_log.exists():
        train_status = "COMPLETED" if parse_exit_code(train_log) == 0 else "FAILED"
        entries.append(
            PipelineEntry(
                timestamp=fmt_time(coerce_dt(train_log, cursor)),
                step="TRAIN",
                status=train_status,
                detail=f"log={train_log} adapter_dir={adapter_dir}",
            )
        )
        cursor = coerce_dt(train_log, cursor)
    else:
        if route_type == "baseline_infer":
            detail = "baseline inference route reuses an existing model"
        elif preview_meta.get("preview_only"):
            detail = "preview-only adapter generated plans without executing training"
        else:
            detail = "no train.log detected"
        entries.append(PipelineEntry(fmt_time(cursor), "TRAIN", "SKIPPED", detail))

    export_log = log_dir / "export.log"
    merged_dir = output_dir / "merged"
    cursor = next_time(cursor, now)
    if export_log.exists():
        export_status = "COMPLETED" if parse_exit_code(export_log) == 0 else "FAILED"
        entries.append(
            PipelineEntry(
                timestamp=fmt_time(coerce_dt(export_log, cursor)),
                step="EXPORT",
                status=export_status,
                detail=f"log={export_log} merged_dir={merged_dir}",
            )
        )
        cursor = coerce_dt(export_log, cursor)
    else:
        if preview_meta.get("preview_only"):
            detail = "preview-only adapter generated export plan only"
        elif route_type == "baseline_infer":
            detail = "baseline inference route does not execute export"
        else:
            detail = "no export.log detected"
        entries.append(PipelineEntry(fmt_time(cursor), "EXPORT", "SKIPPED", detail))

    local_gguf = output_dir / "model.gguf"
    cursor = next_time(cursor, now)
    if local_gguf.exists():
        entries.append(PipelineEntry(fmt_time(coerce_dt(local_gguf, cursor)), "GGUF", "COMPLETED", f"path={local_gguf}"))
        cursor = coerce_dt(local_gguf, cursor)
    elif merged_dir.exists():
        entries.append(
            PipelineEntry(
                fmt_time(cursor),
                "GGUF",
                "PLANNED",
                f"merged model is ready; current phase still needs an explicit GGUF conversion step. merged_dir={merged_dir}",
            )
        )
    else:
        entries.append(PipelineEntry(fmt_time(cursor), "GGUF", "SKIPPED", "no merged model or GGUF artifact detected"))

    cursor = next_time(cursor, now)
    if merged_dir.exists():
        entries.append(
            PipelineEntry(
                fmt_time(cursor),
                "VLLM_REGISTER",
                "PLANNED",
                f"merged model can be served by vLLM directly. model_dir={merged_dir}",
            )
        )
    else:
        entries.append(PipelineEntry(fmt_time(cursor), "VLLM_REGISTER", "SKIPPED", "no merged model directory available"))

    ollama_log = log_dir / "ollama-create.log"
    cursor = next_time(cursor, now)
    if ollama_log.exists():
        ollama_status = "COMPLETED" if parse_exit_code(ollama_log) == 0 else "FAILED"
        first_line = safe_read_text(ollama_log).splitlines()[0] if safe_read_text(ollama_log).splitlines() else "ollama create"
        entries.append(
            PipelineEntry(
                fmt_time(coerce_dt(ollama_log, cursor)),
                "OLLAMA_REGISTER",
                ollama_status,
                first_line,
            )
        )
        cursor = coerce_dt(ollama_log, cursor)
    elif local_gguf.exists():
        entries.append(PipelineEntry(fmt_time(cursor), "OLLAMA_REGISTER", "PLANNED", "GGUF exists but no create log was found yet"))
    else:
        entries.append(
            PipelineEntry(
                fmt_time(cursor),
                "OLLAMA_REGISTER",
                "SKIPPED",
                "no ollama-create.log detected; current phase can backfill after GGUF is ready",
            )
        )

    eval_summary_path = eval_dir / "eval_summary.json"
    cursor = next_time(cursor, now)
    if eval_summary:
        entries.append(
            PipelineEntry(
                fmt_time(coerce_dt(eval_summary_path, cursor)),
                "EVAL",
                "COMPLETED",
                (
                    f"case_count={eval_summary.get('case_count')} "
                    f"contains_accuracy={eval_summary.get('contains_accuracy')} "
                    f"json_valid_rate={eval_summary.get('json_valid_rate')}"
                ),
            )
        )
        cursor = coerce_dt(eval_summary_path, cursor)
    else:
        entries.append(PipelineEntry(fmt_time(cursor), "EVAL", "SKIPPED", "no eval_summary.json detected"))

    benchmark_summary_path = benchmark_dir / "benchmark_summary.json"
    cursor = next_time(cursor, now)
    if benchmark_summary:
        entries.append(
            PipelineEntry(
                fmt_time(coerce_dt(benchmark_summary_path, cursor)),
                "BENCHMARK",
                "COMPLETED",
                (
                    f"case_count={benchmark_summary.get('case_count')} "
                    f"avg_total_ms={benchmark_summary.get('avg_total_ms')} "
                    f"avg_eval_tokens_per_sec={benchmark_summary.get('avg_eval_tokens_per_sec')}"
                ),
            )
        )
        cursor = coerce_dt(benchmark_summary_path, cursor)
    else:
        entries.append(PipelineEntry(fmt_time(cursor), "BENCHMARK", "SKIPPED", "no benchmark_summary.json detected"))

    cursor = next_time(cursor, now)
    if report_path.exists():
        entries.append(PipelineEntry(fmt_time(coerce_dt(report_path, cursor)), "REPORT", "COMPLETED", f"report_path={report_path}"))
        cursor = coerce_dt(report_path, cursor)
        entries.append(PipelineEntry(fmt_time(next_time(cursor, now)), "RUN", "COMPLETED", f"run_id={run_dir.name}"))
    else:
        entries.append(PipelineEntry(fmt_time(cursor), "REPORT", "SKIPPED", "no single run report detected"))

    return entries


def write_pipeline_files(run_dir: Path, entries: list[PipelineEntry]) -> tuple[Path, Path]:
    log_dir = run_dir / "logs"
    log_dir.mkdir(parents=True, exist_ok=True)
    log_path = log_dir / "automation_pipeline.log"
    summary_path = log_dir / "automation_pipeline.json"

    lines = [f"{entry.timestamp} [{entry.status}] {entry.step}: {entry.detail}" for entry in entries]
    log_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    payload = {
        "steps": [
            {
                "timestamp": entry.timestamp,
                "step": entry.step,
                "status": entry.status,
                "detail": entry.detail,
            }
            for entry in entries
        ]
    }
    summary_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return log_path, summary_path


def load_font(size: int, *, mono: bool = False, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    font_candidates = []
    if mono:
        font_candidates.extend(
            [
                Path(r"C:\Windows\Fonts\consola.ttf"),
                Path(r"C:\Windows\Fonts\cour.ttf"),
            ]
        )
    elif bold:
        font_candidates.extend(
            [
                Path(r"C:\Windows\Fonts\msyhbd.ttc"),
                Path(r"C:\Windows\Fonts\simhei.ttf"),
            ]
        )
    else:
        font_candidates.extend(
            [
                Path(r"C:\Windows\Fonts\msyh.ttc"),
                Path(r"C:\Windows\Fonts\simhei.ttf"),
            ]
        )

    for path in font_candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    if not text:
        return [""]
    lines: list[str] = []
    current = ""
    for char in text:
        trial = current + char
        bbox = draw.textbbox((0, 0), trial, font=font)
        if bbox[2] - bbox[0] <= max_width or not current:
            current = trial
            continue
        lines.append(current)
        current = char
    if current:
        lines.append(current)
    return lines


def render_pipeline_image(run_dir: Path, entries: list[PipelineEntry], output_path: Path) -> Path:
    width, height = 1720, 980
    bg = "#0B1020"
    panel = "#11182B"
    border = "#23314E"
    title_color = "#F7FAFC"
    subtitle_color = "#94A3B8"
    detail_color = "#CBD5E1"
    grid_color = "#1E293B"

    image = Image.new("RGB", (width, height), bg)
    draw = ImageDraw.Draw(image)
    title_font = load_font(34, bold=True)
    subtitle_font = load_font(18)
    body_font = load_font(24, mono=True)
    detail_font = load_font(19)
    tag_font = load_font(18, bold=True)

    outer = (54, 58, width - 54, height - 58)
    draw.rounded_rectangle(outer, radius=24, fill=panel, outline=border, width=2)
    header_h = 88
    draw.rounded_rectangle((outer[0], outer[1], outer[2], outer[1] + header_h), radius=24, fill="#141F36", outline=border, width=2)
    draw.rectangle((outer[0], outer[1] + 24, outer[2], outer[1] + header_h), fill="#141F36")

    for idx, color in enumerate(["#FF5F56", "#FFBD2E", "#27C93F"]):
        cx = outer[0] + 34 + idx * 26
        cy = outer[1] + 28
        draw.ellipse((cx - 7, cy - 7, cx + 7, cy + 7), fill=color)

    draw.text((outer[0] + 98, outer[1] + 14), "模型训练工作台 - 自动化流程日志", fill=title_color, font=title_font)
    subtitle = (
        f"运行目录: {run_dir.name}    "
        "流程: export -> GGUF -> Ollama/vLLM 注册 -> 评测    "
        "说明: 当前截图按现有运行产物回填生成"
    )
    draw.text((outer[0] + 98, outer[1] + 54), subtitle, fill=subtitle_color, font=subtitle_font)

    top = outer[1] + header_h + 18
    left = outer[0] + 26
    right = outer[2] - 26
    line_y = top

    for entry in entries:
        draw.line((left, line_y + 66, right, line_y + 66), fill=grid_color, width=1)
        status_color = STATUS_COLORS.get(entry.status, "#A0AEC0")
        tag_box = (left + 120, line_y + 12, left + 248, line_y + 48)
        draw.rounded_rectangle(tag_box, radius=10, fill=status_color)
        draw.text((left + 141, line_y + 18), STATUS_LABELS.get(entry.status, entry.status), fill="#081018", font=tag_font)

        draw.text((left + 8, line_y + 18), entry.timestamp[-8:], fill="#67E8F9", font=body_font)
        step_label = STEP_LABELS.get(entry.step, entry.step)
        draw.text((left + 278, line_y + 16), step_label, fill=title_color, font=body_font)

        wrapped = wrap_text(draw, entry.detail, detail_font, max_width=right - (left + 430))
        detail_y = line_y + 18
        for idx, segment in enumerate(wrapped[:2]):
            draw.text((left + 430, detail_y + idx * 24), segment, fill=detail_color, font=detail_font)

        line_y += 76
        if line_y > outer[3] - 120:
            break

    footer = "注: 一期已支持统一流水日志落盘; 若未实际执行 GGUF/vLLM, 日志会明确标注为预留或跳过。"
    draw.text((outer[0] + 28, outer[3] - 40), footer, fill=subtitle_color, font=subtitle_font)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path)
    return output_path


def insert_image_into_ppt(ppt_path: Path, image_path: Path) -> Path:
    prs = Presentation(str(ppt_path))
    slide = prs.slides[15]
    placeholder = slide.shapes[5]
    slide.shapes.add_picture(
        str(image_path),
        placeholder.left,
        placeholder.top,
        width=placeholder.width,
        height=placeholder.height,
    )
    temp_output = ppt_path.with_name(f"{ppt_path.stem}.tmp_patch{ppt_path.suffix}")
    prs.save(str(temp_output))
    try:
        temp_output.replace(ppt_path)
        return ppt_path
    except PermissionError:
        fallback = ppt_path.with_name(f"{ppt_path.stem}-自动化流程日志版{ppt_path.suffix}")
        temp_output.replace(fallback)
        return fallback


def main() -> None:
    parser = argparse.ArgumentParser(description="Backfill automation pipeline assets and inject a screenshot into the PPT.")
    parser.add_argument(
        "--run-dir",
        default=r"E:\ai\ai_trains\runtime\runs\run_20260412_053225_425ae0",
        help="Run directory used to synthesize the pipeline log and screenshot.",
    )
    parser.add_argument(
        "--image-output",
        default=r"E:\ai\ai_trains\resources\screenshots\automation_pipeline_log.png",
        help="PNG output path for the generated automation pipeline screenshot.",
    )
    parser.add_argument(
        "--resources-dir",
        default=r"E:\ai\ai_trains\resources",
        help="Resources directory that contains the training introduction PPT.",
    )
    parser.add_argument(
        "--skip-ppt",
        action="store_true",
        help="Generate the pipeline log and screenshot without modifying the PPT.",
    )
    args = parser.parse_args()

    run_dir = Path(args.run_dir)
    image_output = Path(args.image_output)
    resources_dir = Path(args.resources_dir)

    entries = build_entries_for_run(run_dir)
    write_pipeline_files(run_dir, entries)
    render_pipeline_image(run_dir, entries, image_output)

    if not args.skip_ppt:
        ppt_path = find_training_intro_ppt(resources_dir)
        final_ppt = insert_image_into_ppt(ppt_path, image_output)
        print(final_ppt)

    print(image_output)


if __name__ == "__main__":
    main()
