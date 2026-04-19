from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
MID_BLUE = RGBColor(0x2B, 0x57, 0x9A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
CODE_WHITE = RGBColor(0xF8, 0xF8, 0xF2)
CODE_GREEN = RGBColor(0xA6, 0xE2, 0x2E)
CODE_YELLOW = RGBColor(0xE6, 0xDB, 0x74)
CODE_COMMENT = RGBColor(0x75, 0x71, 0x5E)

REPO_ROOT = Path(__file__).resolve().parents[1]


def find_target_ppt() -> Path:
    resources = REPO_ROOT / "resources"
    return next(
        p
        for p in resources.glob("*.pptx")
        if "Qwen3.5-9B-Base" in p.name and "bak_fillshots" not in p.name and not p.name.startswith("~$")
    )


def clear_and_set_lines(
    shape,
    lines: list[str],
    *,
    font_name: str = "Microsoft YaHei",
    font_size: int = 18,
    color: RGBColor = DARK_GRAY,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = alignment
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color


def set_bullets(shape, items: list[str], *, font_size: int = 14, color: RGBColor = DARK_GRAY) -> None:
    clear_and_set_lines(
        shape,
        [f"• {item}" for item in items],
        font_name="Microsoft YaHei",
        font_size=font_size,
        color=color,
        bold=False,
        alignment=PP_ALIGN.LEFT,
    )


def set_code(shape, lines: list[tuple[str, RGBColor]]) -> None:
    tf = shape.text_frame
    tf.clear()
    for idx, (line, color) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.bold = False
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT


def set_placeholder(shape, title: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = f"截图预留：{title}"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "Microsoft YaHei"
    p1.font.size = Pt(16)
    p1.font.color.rgb = MID_GRAY

    p2 = tf.add_paragraph()
    p2.text = "（请在此处粘贴实际截图）"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(12)
    p2.font.color.rgb = MID_GRAY


def set_table(shape, rows: list[list[str]]) -> None:
    table = shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(11)
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.bold = False
                    p.font.color.rgb = DARK_GRAY


def repair_slide_1(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["基于 Qwen3.5-9B-Base 的动环监控"], font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    clear_and_set_lines(slide.shapes[1], ["告警诊断场景 LoRA-SFT 训练介绍"], font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    clear_and_set_lines(slide.shapes[2], ["训练框架：ms-swift  |  训练方式：LoRA-SFT  |  重点案例：温度过高告警"], font_size=16, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
    clear_and_set_lines(slide.shapes[3], ["ai_trains 训练工作台"], font_size=14, color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)


def repair_slide_2(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["目录"], font_size=32, bold=True, color=DARK_BLUE)
    set_bullets(
        slide.shapes[1],
        [
            "一、项目背景与业务场景",
            "二、模型选型：Qwen3.5-9B-Base",
            "三、数据准备：温度过高告警 SFT 数据",
            "四、环境搭建：ms-swift 安装配置",
            "五、训练配置：LoRA-SFT 参数",
            "六、训练执行：swift sft 命令",
            "七、模型导出：LoRA 合并与量化",
            "八、模型部署：Ollama 与 vLLM",
            "九、推理验证：vLLM 温度过高告警案例",
            "十、效果评测：离线评测与对比报告",
            "十一、平台集成：训练工作台闭环",
            "十二、运行产物与效果展示",
        ],
        font_size=16,
    )


def repair_slide_3(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["一、项目背景与业务场景"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["动环监控平台 × 大模型告警诊断"], font_size=18, bold=True, color=MID_BLUE)
    set_bullets(
        slide.shapes[2],
        [
            "业务痛点：告警量大、跨设备关联复杂，人工排障难以稳定输出根因、影响范围和处置建议",
            "核心目标：围绕模型训练构建可复用工作台，沉淀训练、推理验证、对比评测三类能力",
            "专业场景：告警诊断只是首个方向，后续可扩展巡检问答、知识库问答、运维工单生成等场景",
            "示例主线：本文以“温湿度01 温度高”告警为例，贯穿训练样本、推理请求和 vLLM 验证",
            "目标模型：优先使用 Qwen3.5-9B-Base + ms-swift LoRA-SFT，兼容一期 LLaMA-Factory 方案",
        ],
        font_size=13,
    )
    set_placeholder(slide.shapes[3], "动环监控平台告警页面截图")
    set_placeholder(slide.shapes[4], "温度过高告警诊断结果页截图")


def repair_slide_4(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["一、业务场景：温度过高告警训练/推理样例"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["统一用一个案例讲清训练输入、推理输出和评测口径"], font_size=18, bold=True, color=MID_BLUE)
    set_code(
        slide.shapes[2],
        [
            ("// 输入（温度过高告警样例）", CODE_COMMENT),
            ("告警设备名称：温湿度01", CODE_WHITE),
            ("告警设备编码：03202206221157_170100000000003", CODE_YELLOW),
            ("告警开始时间：2025-12-18 19:28:48", CODE_WHITE),
            ("告警信号名称：温度高", CODE_GREEN),
            ("测点：017301000 / 温度01 / 35.2度", CODE_WHITE),
            ("采集时间：2023-10-05 14:30:00", CODE_WHITE),
            ("伴随现象：门磁告警同时出现", CODE_YELLOW),
            ("同站点其他区域温度正常", CODE_WHITE),
            ("空调运行状态显示异常", CODE_GREEN),
            ("候选方向：空调系统故障 / 温度传感器故障", CODE_WHITE),
        ],
    )
    set_code(
        slide.shapes[3],
        [
            ("// 输出（期望诊断报告）", CODE_COMMENT),
            ("诊断结论：机房温度异常升高，", CODE_GREEN),
            ("主要由空调系统故障引起。", CODE_GREEN),
            ("处理建议：", CODE_YELLOW),
            ("1. 优先排查空调压缩机运行状态", CODE_WHITE),
            ("2. 查看空调系统日志定位异常原因", CODE_WHITE),
            ("3. 校准温度传感器，排除误报", CODE_WHITE),
            ("成因分析：", CODE_YELLOW),
            ("1. 空调系统故障导致制冷不足 (85%)", CODE_WHITE),
            ("   证据：温度持续 >35℃，空调状态异常", CODE_WHITE),
            ("2. 温度传感器故障 (70%)", CODE_WHITE),
            ("   证据：其他传感器正常，需复核温感", CODE_WHITE),
        ],
    )
    clear_and_set_lines(
        slide.shapes[4],
        ["分析过程建议保留五步：检查告警时间点 → 对比同站点数据 → 分析历史趋势 → 排除传感器故障 → 核验空调日志"],
        font_size=12,
        bold=True,
        color=MID_BLUE,
    )


def repair_slide_5(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["二、模型选型：Qwen3.5-9B-Base"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["为什么在本项目中优先选择 9B-Base"], font_size=18, bold=True, color=MID_BLUE)
    set_table(
        slide.shapes[2],
        [
            ["维度", "说明"],
            ["模型规模", "9B 参数量，适合作为动环垂直领域训练的主力模型"],
            ["Base vs Instruct", "Base 模型便于通过 SFT 注入统一的诊断格式和专业术语风格"],
            ["中文与工程能力", "适配设备名、测点值、时间戳等中文结构化文本"],
            ["长上下文", "可容纳较长告警摘要、历史日志与多设备关联信息"],
            ["部署兼容性", "训练后既可导出 GGUF 给 Ollama，也可直接使用 merged 模型给 vLLM"],
            ["二期扩展性", "后续可平滑扩展到 Qwen3.5-9B 与其他专业场景"],
        ],
    )
    clear_and_set_lines(
        slide.shapes[3],
        ["模型 ID：Qwen/Qwen3.5-9B-Base    模板：qwen3    训练框架：ms-swift (LoRA-SFT)"],
        font_size=14,
        bold=True,
        color=MID_BLUE,
    )


def repair_slide_6(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["三、数据准备：温度过高告警 SFT 数据"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["训练集来自真实告警，测试集围绕高温场景做模拟补充"], font_size=18, bold=True, color=MID_BLUE)
    clear_and_set_lines(
        slide.shapes[2],
        ["数据加工链路：原始告警 JSON → 压缩整理 → 根因/建议抽取 → SFT 样本构建 → 训练/评测集切分"],
        font_size=13,
        bold=True,
        color=MID_BLUE,
    )
    set_code(
        slide.shapes[3],
        [
            ("# 生成训练/评测集", CODE_COMMENT),
            ("python scripts/prepare_alarm_analysis_datasets.py \\", CODE_GREEN),
            ("  --source resources/alarm_analysis.json \\", CODE_WHITE),
            ("  --datasets-root runtime/datasets \\", CODE_WHITE),
            ("  --swift-train-version alarm_analysis_swift_sft_v1 \\", CODE_YELLOW),
            ("  --eval-version alarm_analysis_eval_v1 \\", CODE_WHITE),
            ("  --eval-count 64 --seed 42", CODE_WHITE),
        ],
    )
    set_code(
        slide.shapes[4],
        [
            ("{", CODE_WHITE),
            ('  "messages": [', CODE_WHITE),
            ('    {"role":"system","content":"你是机房设备运维专家"},', CODE_YELLOW),
            ('    {"role":"user","content":"告警设备名称：温湿度01 ..."},', CODE_GREEN),
            ('    {"role":"assistant","content":"诊断结论：机房温度异常升高 ..."}', CODE_YELLOW),
            ("  ]", CODE_WHITE),
            ("}", CODE_WHITE),
        ],
    )
    set_table(
        slide.shapes[5],
        [
            ["数据集/子集", "格式", "用途", "说明"],
            ["alarm_analysis_swift_sft_v1", "JSONL(messages)", "ms-swift 训练", "主训练集，覆盖空调/UPS/市电/温湿度等场景"],
            ["temperature_high_train_cases", "样例子集", "专题抽样检查", "从真实样本中抽取温度高/温度过高案例用于人工复核"],
            ["temperature_high_eval_mock", "JSONL", "推理验证", "模拟 6-12 条高温样例，覆盖空调故障和传感器误报两类标签"],
        ],
    )


def repair_slide_7(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["四、环境搭建：ms-swift 安装配置"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["建议在 WSL 中部署，根目录按项目约定放在 E:\\wsl"], font_size=18, bold=True, color=MID_BLUE)
    set_code(
        slide.shapes[2],
        [
            ("# 1. 创建 conda 环境", CODE_COMMENT),
            ("conda create -n swift python=3.11 -y", CODE_GREEN),
            ("conda activate swift", CODE_GREEN),
            ("", CODE_WHITE),
            ("# 2. 安装 PyTorch (CUDA 12.1)", CODE_COMMENT),
            ("pip install torch torchvision torchaudio \\", CODE_GREEN),
            ("  --index-url https://download.pytorch.org/whl/cu121", CODE_WHITE),
            ("", CODE_WHITE),
            ("# 3. 安装 ms-swift", CODE_COMMENT),
            ("pip install 'ms-swift[llm]' -U", CODE_GREEN),
            ("", CODE_WHITE),
            ("# 4. 验证", CODE_COMMENT),
            ("swift sft --help", CODE_GREEN),
        ],
    )
    set_code(
        slide.shapes[3],
        [
            ("# 5. 下载基础模型", CODE_COMMENT),
            ("export USE_MODELSCOPE_HUB=1", CODE_GREEN),
            ("modelscope download \\", CODE_GREEN),
            ("  --model Qwen/Qwen3.5-9B-Base \\", CODE_WHITE),
            ("  --local_dir /data/models/Qwen3.5-9B-Base", CODE_WHITE),
            ("", CODE_WHITE),
            ("# 6. 工作目录", CODE_COMMENT),
            ("mkdir -p ./runtime/datasets", CODE_GREEN),
            ("mkdir -p ./runtime/runs", CODE_GREEN),
            ("", CODE_WHITE),
            ("# 7. GPU 检查", CODE_COMMENT),
            ('python -c "import torch; print(torch.cuda.is_available())"', CODE_YELLOW),
        ],
    )
    clear_and_set_lines(slide.shapes[4], ["硬件建议：单卡 24GB+ 可做 LoRA-SFT；若要更稳妥跑 9B，建议 4090 / A100 级别显卡"], font_size=13, bold=True, color=RGBColor(0xE8, 0x6C, 0x00))
    clear_and_set_lines(slide.shapes[5], ["软件建议：Python 3.11 + PyTorch 2.x + CUDA 12.1+ + ms-swift 最新版本"], font_size=13, bold=True, color=MID_BLUE)


def repair_slide_8(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["五、训练配置：LoRA-SFT 参数"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["重点关注训练参数如何影响告警诊断质量"], font_size=18, bold=True, color=MID_BLUE)
    set_code(
        slide.shapes[2],
        [
            ("# swift.train.yaml", CODE_COMMENT),
            ("model: Qwen/Qwen3.5-9B-Base", CODE_GREEN),
            ("template: qwen3", CODE_WHITE),
            ("dataset:", CODE_YELLOW),
            ("  - runtime/datasets/alarm_analysis_swift_sft_v1/", CODE_WHITE),
            ("    alarm_analysis_swift_sft_v1.jsonl", CODE_WHITE),
            ("train_type: lora", CODE_GREEN),
            ("target_modules: all-linear", CODE_WHITE),
            ("lora_rank: 64", CODE_YELLOW),
            ("lora_alpha: 128", CODE_YELLOW),
            ("max_length: 4096", CODE_WHITE),
            ("per_device_train_batch_size: 1", CODE_WHITE),
            ("gradient_accumulation_steps: 8", CODE_WHITE),
            ("learning_rate: 2e-4", CODE_GREEN),
            ("num_train_epochs: 3", CODE_WHITE),
            ("output_dir: runtime/runs/{run_id}/output/adapter", CODE_GREEN),
        ],
    )
    set_table(
        slide.shapes[3],
        [
            ["参数", "建议值", "原因"],
            ["template", "qwen3", "与 Qwen3.5-9B-Base 对齐"],
            ["lora_rank", "64", "兼顾领域表达能力与显存占用"],
            ["lora_alpha", "128", "常见稳妥组合"],
            ["target_modules", "all-linear", "充分学习诊断结构和术语"],
            ["max_length", "4096", "容纳较长告警摘要与分析过程"],
            ["learning_rate", "2e-4", "LoRA-SFT 常见起点"],
            ["epochs", "3", "一期先小步快跑，避免过拟合"],
            ["grad_accum", "8", "单卡下提升等效批量"],
        ],
    )


def repair_slide_9(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["六、训练执行：swift sft 命令"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["直接落到命令行，便于后端 Worker 编排"], font_size=18, bold=True, color=MID_BLUE)
    clear_and_set_lines(slide.shapes[2], ["方式一：使用配置文件（推荐）"], font_size=14, bold=True, color=DARK_GRAY)
    set_code(slide.shapes[3], [("CUDA_VISIBLE_DEVICES=0 swift sft --config swift.train.yaml", CODE_GREEN)])
    clear_and_set_lines(slide.shapes[4], ["方式二：命令行直接指定参数"], font_size=14, bold=True, color=DARK_GRAY)
    set_code(
        slide.shapes[5],
        [
            ("CUDA_VISIBLE_DEVICES=0 swift sft \\", CODE_GREEN),
            ("  --model Qwen/Qwen3.5-9B-Base \\", CODE_WHITE),
            ("  --dataset runtime/datasets/alarm_analysis_swift_sft_v1/alarm_analysis_swift_sft_v1.jsonl \\", CODE_WHITE),
            ("  --train_type lora --template qwen3 \\", CODE_YELLOW),
            ("  --max_length 4096 --num_train_epochs 3 \\", CODE_WHITE),
            ("  --per_device_train_batch_size 1 \\", CODE_WHITE),
            ("  --gradient_accumulation_steps 8 \\", CODE_WHITE),
            ("  --lora_rank 64 --lora_alpha 128 \\", CODE_YELLOW),
            ("  --target_modules all-linear \\", CODE_WHITE),
            ("  --output_dir runtime/runs/{run_id}/output/adapter", CODE_GREEN),
        ],
    )
    clear_and_set_lines(slide.shapes[6], ["后端编排建议：Worker 负责生成 run_id、写入 swift.train.yaml、启动训练并采集 train.log。"], font_size=12, bold=True, color=MID_BLUE)


def repair_slide_10(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["七、模型导出：LoRA 合并与量化"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["导出阶段需要同时照顾 vLLM 与 Ollama 两条推理链路"], font_size=18, bold=True, color=MID_BLUE)
    set_code(
        slide.shapes[2],
        [
            ("# 步骤一：合并 LoRA 适配器到基础模型", CODE_COMMENT),
            ("swift export \\", CODE_GREEN),
            ("  --model Qwen/Qwen3.5-9B-Base \\", CODE_WHITE),
            ("  --adapters ./runtime/runs/{run_id}/output/adapter \\", CODE_YELLOW),
            ("  --merge_lora true \\", CODE_WHITE),
            ("  --output_dir ./runtime/runs/{run_id}/output/merged", CODE_GREEN),
        ],
    )
    set_code(
        slide.shapes[3],
        [
            ("# 步骤二：量化为 GGUF（仅 Ollama 需要）", CODE_COMMENT),
            ("python llama.cpp/convert_hf_to_gguf.py \\", CODE_GREEN),
            ("  ./runtime/runs/{run_id}/output/merged \\", CODE_WHITE),
            ("  --outfile ./runtime/artifacts/gguf/alarm-qwen35-9b-sft.gguf \\", CODE_YELLOW),
            ("  --outtype q4_k_m", CODE_GREEN),
        ],
    )
    clear_and_set_lines(slide.shapes[4], ["关键结论：vLLM 优先直接读取 merged/ 完整 Hugging Face 模型目录；GGUF 只用于 Ollama / llama.cpp 对比推理。"], font_size=13, bold=True, color=RGBColor(0xE8, 0x6C, 0x00))
    set_placeholder(slide.shapes[5], "导出命令执行结果截图")


def repair_slide_11(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["八、模型部署：Ollama 与 vLLM"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["一期保留 Ollama 对比，生产/正式验证优先走 vLLM"], font_size=18, bold=True, color=MID_BLUE)
    set_code(
        slide.shapes[2],
        [
            ("# Ollama：本地对比推理链路", CODE_COMMENT),
            ("FROM ./alarm-qwen35-9b-sft.gguf", CODE_GREEN),
            ("PARAMETER temperature 0.1", CODE_WHITE),
            ("PARAMETER top_p 0.8", CODE_WHITE),
            ("PARAMETER num_ctx 8192", CODE_GREEN),
            ("", CODE_WHITE),
            ("ollama create alarm-swift-qwen35-9b -f Modelfile", CODE_GREEN),
            ("ollama run alarm-swift-qwen35-9b", CODE_YELLOW),
        ],
    )
    set_code(
        slide.shapes[3],
        [
            ("# vLLM：推荐部署链路（读取 merged 模型）", CODE_COMMENT),
            ("vllm serve ./runtime/runs/{run_id}/output/merged \\", CODE_GREEN),
            ("  --served-model-name alarm-qwen35-9b-sft \\", CODE_WHITE),
            ("  --host 0.0.0.0 --port 8000 \\", CODE_WHITE),
            ("  --dtype auto \\", CODE_WHITE),
            ("  --api-key train-workbench \\", CODE_YELLOW),
            ("  --generation-config vllm \\", CODE_WHITE),
            ("  --language-model-only", CODE_GREEN),
        ],
    )
    set_bullets(
        slide.shapes[4],
        [
            "Ollama 适合一期平台统一本地对比评测，接入简单，便于桌面化验证。",
            "vLLM 适合正式部署与 OpenAI 兼容接口验证，可直接对接后端服务或工作台推理模块。",
            "同一轮训练产物建议同时保留 merged/ 与 .gguf，分别服务 vLLM 与 Ollama。",
        ],
        font_size=12,
    )


def repair_slide_12(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["九、推理验证：vLLM 回风温度过高告警案例"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["验证服务可用性、接口正确性和诊断结果完整性"], font_size=18, bold=True, color=MID_BLUE)
    set_code(
        slide.shapes[2],
        [
            ("# 1. 检查模型是否已注册", CODE_COMMENT),
            ('curl http://127.0.0.1:8000/v1/models \\', CODE_GREEN),
            ('  -H "Authorization: Bearer train-workbench"', CODE_WHITE),
            ("", CODE_WHITE),
            ("# 2. 发起回风温度过高告警诊断请求", CODE_COMMENT),
            ("curl http://127.0.0.1:8000/v1/chat/completions \\", CODE_GREEN),
            ('  -H "Content-Type: application/json" \\', CODE_WHITE),
            ('  -H "Authorization: Bearer train-workbench" \\', CODE_WHITE),
            ('  -d \'{"model":"alarm-qwen35-9b-sft",', CODE_YELLOW),
            ('        "messages":[{"role":"system","content":"你是机房设备运维专家"},', CODE_WHITE),
            ('        {"role":"user","content":"告警设备名称：普通空调-1#空调内机-科龙\\n告警信号名称：回风温度过高告警\\n回风温度1：33.00℃\\n设定温度：25.00℃..."}],', CODE_WHITE),
            ('        "temperature":0.1,"max_tokens":512}\'', CODE_YELLOW),
        ],
    )
    set_code(
        slide.shapes[3],
        [
            ("# 期望响应片段（示例）", CODE_COMMENT),
            ("诊断结论：回风温度过高告警主要由", CODE_GREEN),
            ("空调制冷系统故障导致（85%可能性）。", CODE_GREEN),
            ("", CODE_WHITE),
            ("处理建议：", CODE_YELLOW),
            ("1. 优先检查压缩机、冷媒压力和过滤网", CODE_WHITE),
            ("2. 测量出风/回风温差验证制冷效果", CODE_WHITE),
            ("3. 复测回风温度传感器并检查线路", CODE_WHITE),
            ("", CODE_WHITE),
            ("成因分析：", CODE_YELLOW),
            ("1. 空调制冷系统故障 (85%)", CODE_WHITE),
            ("2. 温度传感器故障 (15%)", CODE_WHITE),
        ],
    )
    clear_and_set_lines(
        slide.shapes[4],
        ["验证通过标准：接口返回正常、输出包含主因/次因、处置建议，并引用关键证据（回风温度33.00℃/设定温度25.00℃/环境温度26.30-26.60℃/制冷模式运行）。"],
        font_size=12,
        bold=True,
        color=MID_BLUE,
    )


def repair_slide_13(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["十、效果评测：离线评测与对比报告"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["围绕回风温度过高告警建立明确的评测口径"], font_size=18, bold=True, color=MID_BLUE)
    set_bullets(
        slide.shapes[2],
        [
            "评测集：alarm_analysis_eval_v1 + return_air_temp_eval_mock",
            "Smoke 集：先抽 2-4 条回风温度过高样例做快速验证",
            "重点指标：主因命中、次因命中、处置建议完整度、关键证据引用",
            "评测示例：普通空调-1#空调内机-科龙 / 回风温度过高告警 / 2026-03-10 10:23:45",
        ],
        font_size=12,
    )
    set_code(
        slide.shapes[3],
        [
            ("# 评测样例（回风温度过高告警）", CODE_COMMENT),
            ("{", CODE_WHITE),
            ('  "task_type": "alarm_diagnosis",', CODE_YELLOW),
            ('  "messages": [', CODE_WHITE),
            ('    "告警设备名称：普通空调-1#空调内机-科龙",', CODE_GREEN),
            ('    "告警设备编码：02202205110102_150100000000001",', CODE_WHITE),
            ('    "告警开始时间：2026-03-10 10:23:45",', CODE_WHITE),
            ('    "告警信号名称：回风温度过高告警",', CODE_YELLOW),
            ('    "回风温度1=33.00℃；设定温度=25.00℃；环境温度=26.30-26.60℃",', CODE_WHITE),
            ('    "工作状态=开机；运行模式=制冷模式；直流输出电压=53.50-54.00V"', CODE_WHITE),
            ('  ],', CODE_WHITE),
            ('  "expected_contains": [', CODE_GREEN),
            ('    "空调制冷系统故障",', CODE_WHITE),
            ('    "温度传感器故障",', CODE_WHITE),
            ('    "压缩机运行状态"', CODE_WHITE),
            ("  ]", CODE_WHITE),
            ("}", CODE_WHITE),
        ],
    )
    set_table(
        slide.shapes[4],
        [
            ["验证维度", "Qwen3.5-9B-Base 基线", "LoRA-SFT (swift) 目标"],
            ["根因定位", "能否识别空调高温异常", "稳定输出主因=空调制冷系统故障，次因=温度传感器故障"],
            ["处置建议", "建议偏泛化", "覆盖压缩机/冷媒/过滤网检查 + 传感器复测与线路核验"],
            ["证据引用", "不一定引用测点和数值", "引用回风温度33.00℃、设定温度25.00℃、环境温度26.30-26.60℃、制冷模式运行"],
            ["接口验证", "Ollama 对比输出", "vLLM OpenAI 兼容接口可稳定返回"],
        ],
    )
    set_placeholder(slide.shapes[5], "回风温度过高评测结果对比图表截图")
    set_placeholder(slide.shapes[6], "回风温度过高对比报告页面截图")


def repair_slide_14(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["十一、平台集成：训练工作台全链路闭环"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["训练、导出、推理验证、对比评测在同一工作台收口"], font_size=18, bold=True, color=MID_BLUE)
    set_bullets(
        slide.shapes[2],
        [
            "前端：Vue3 + Vite H5 训练工作台",
            "后端：FastAPI + SQLite 元数据管理",
            "Worker：串行执行训练、导出、评测流程",
            "训练后端：LLaMA-Factory（一期） / ms-swift（二期）",
            "推理后端：Ollama（本地对比） / vLLM（正式验证）",
            "场景扩展：告警诊断只是第一批，后续可扩展到其他专业场景与模型",
        ],
        font_size=12,
    )
    set_code(
        slide.shapes[3],
        [
            ("POST /api/experiments", CODE_COMMENT),
            ("{", CODE_WHITE),
            ('  "name": "Qwen3.5-9B-Base 告警诊断训练",', CODE_YELLOW),
            ('  "base_model": "Qwen/Qwen3.5-9B-Base",', CODE_GREEN),
            ('  "trainer_backend": "swift",', CODE_GREEN),
            ('  "route_type": "sft",', CODE_WHITE),
            ('  "dataset_version": "alarm_analysis_swift_sft_v1",', CODE_YELLOW),
            ('  "inference_backend": "vllm",', CODE_YELLOW),
            ('  "train_config": {"train_type":"lora","template":"qwen3"}', CODE_WHITE),
            ("}", CODE_WHITE),
        ],
    )
    set_placeholder(slide.shapes[4], "训练工作台首页截图")
    set_placeholder(slide.shapes[5], "运行详情页（训练/推理）截图")


def repair_slide_15(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["十二、运行产物与目录结构"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["后端训练和部署链路最终沉淀到具体文件产物"], font_size=18, bold=True, color=MID_BLUE)
    set_code(
        slide.shapes[2],
        [
            ("runtime/runs/run_20260412_053225_425ae0/", CODE_GREEN),
            ("├── input/", CODE_WHITE),
            ("│   ├── run_manifest.json", CODE_WHITE),
            ("│   ├── swift.train.yaml      # 训练配置", CODE_YELLOW),
            ("│   ├── swift.export.yaml     # 导出配置", CODE_YELLOW),
            ("│   └── swift.eval.yaml       # 评测计划", CODE_YELLOW),
            ("├── output/", CODE_WHITE),
            ("│   ├── adapter/              # LoRA 权重", CODE_GREEN),
            ("│   ├── merged/               # vLLM 直接读取", CODE_GREEN),
            ("│   ├── Modelfile             # Ollama 对比推理", CODE_WHITE),
            ("│   └── command_preview.json", CODE_WHITE),
            ("├── logs/                     # train/export/worker 日志", CODE_WHITE),
            ("├── eval/                     # 逐条评测结果与汇总", CODE_YELLOW),
            ("└── reports/                  # 运行报告", CODE_WHITE),
        ],
    )
    set_table(
        slide.shapes[3],
        [
            ["产物", "路径", "说明"],
            ["训练配置", "input/swift.train.yaml", "完整训练参数"],
            ["LoRA 权重", "output/adapter/", "训练产出的 adapter"],
            ["合并模型", "output/merged/", "vLLM 部署入口"],
            ["量化模型", "runtime/artifacts/gguf/", "Ollama 对比入口"],
            ["训练日志", "logs/train.log", "训练过程日志"],
            ["评测结果", "eval/eval_results.jsonl", "逐条输出"],
            ["汇总报告", "reports/{run_id}.json", "训练+推理闭环记录"],
        ],
    )


def repair_slide_16(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["十三、效果展示"], font_size=32, bold=True, color=DARK_BLUE)
    clear_and_set_lines(slide.shapes[1], ["预留训练、vLLM 验证和工作台效果截图"], font_size=18, bold=True, color=MID_BLUE)
    set_placeholder(slide.shapes[2], "训练前：Base 模型温度高告警输出截图")
    set_placeholder(slide.shapes[3], "训练后：SFT 模型温度高告警输出截图")
    set_placeholder(slide.shapes[4], "vLLM /v1/chat/completions 验证截图")
    set_placeholder(slide.shapes[5], "训练工作台对比报告截图")


def repair_slide_17(slide) -> None:
    clear_and_set_lines(slide.shapes[0], ["总结与下一步"], font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    set_bullets(
        slide.shapes[1],
        [
            "围绕“温度过高告警”构建统一示例后，训练、推理和评测链路更容易讲透也更便于验收",
            "ms-swift 方案的核心在于：数据格式统一、LoRA-SFT 配置可复用、Worker 易于编排",
            "训练后建议同时保留 merged 与 GGUF，分别服务 vLLM 正式验证和 Ollama 本地对比",
            "vLLM 部署应直接读取 merged 模型目录，并通过 OpenAI 兼容接口完成服务验收",
        ],
        font_size=15,
        color=WHITE,
    )
    clear_and_set_lines(slide.shapes[2], ["下一步建议："], font_size=18, bold=True, color=RGBColor(0xFF, 0xCC, 0x00))
    set_bullets(
        slide.shapes[3],
        [
            "在 WSL 中补齐真实的 Qwen3.5-9B-Base 下载、训练、导出、vLLM 服务验证闭环",
            "补充温度高场景模拟测试集，细化空调故障、传感器误报、门磁伴随告警三类标签",
            "将 vLLM 推理适配进训练工作台后端，形成 Ollama / vLLM 可切换的统一推理层",
        ],
        font_size=14,
        color=RGBColor(0xCC, 0xDD, 0xFF),
    )


def main() -> None:
    target = find_target_ppt()
    prs = Presentation(str(target))

    repair_slide_1(prs.slides[0])
    repair_slide_2(prs.slides[1])
    repair_slide_3(prs.slides[2])
    repair_slide_4(prs.slides[3])
    repair_slide_5(prs.slides[4])
    repair_slide_6(prs.slides[5])
    repair_slide_7(prs.slides[6])
    repair_slide_8(prs.slides[7])
    repair_slide_9(prs.slides[8])
    repair_slide_10(prs.slides[9])
    repair_slide_11(prs.slides[10])
    repair_slide_12(prs.slides[11])
    repair_slide_13(prs.slides[12])
    repair_slide_14(prs.slides[13])
    repair_slide_15(prs.slides[14])
    repair_slide_16(prs.slides[15])
    repair_slide_17(prs.slides[16])

    temp = target.with_name("ppt_training_intro_repaired_temp.pptx")
    prs.save(str(temp))
    print(temp)


if __name__ == "__main__":
    main()
