"""
生成 Qwen3.5-9B-Base 动环监控告警诊断场景训练介绍 PPT。

基于 resources/参考材料补充-训练.pptx 模板风格，
生成一份深入介绍 ms-swift 训练方式的演示文稿。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = REPO_ROOT / "resources" / "参考材料补充-训练.pptx"
OUTPUT_PATH = REPO_ROOT / "resources" / "Qwen3.5-9B-Base动环告警诊断训练介绍.pptx"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
MID_BLUE = RGBColor(0x2B, 0x57, 0x9A)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_GREEN = RGBColor(0xA6, 0xE2, 0x2E)
CODE_YELLOW = RGBColor(0xE6, 0xDB, 0x74)
CODE_WHITE = RGBColor(0xF8, 0xF8, 0xF2)
CODE_COMMENT = RGBColor(0x75, 0x71, 0x5E)


def _add_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                 font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_title(slide, text, font_size=32, color=DARK_BLUE):
    _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 text, font_size=font_size, bold=True, color=color,
                 alignment=PP_ALIGN.LEFT)


def _add_subtitle(slide, text, top=Inches(1.0)):
    _add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.5),
                 text, font_size=18, bold=True, color=MID_BLUE)


def _add_code_block(slide, left, top, width, height, lines: list[tuple[str, RGBColor]]):
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    sp = txBox._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = sp.makeelement(qn("p:spPr"), {})
        sp.append(spPr)
    solidFill = spPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "1E1E2E"})
    solidFill.append(srgbClr)
    spPr.insert(0, solidFill)

    for i, (text, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = color
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    return txBox


def _add_bullet_points(slide, left, top, width, height, items: list[str],
                       font_size=14, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    return txBox


def _add_screenshot_placeholder(slide, left, top, width, height, label="效果截图预留区域"):
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(left, top, width, height)

    sp = txBox._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        spPr = sp.makeelement(qn("p:spPr"), {})
        sp.append(spPr)
    fill_elem = spPr
    solidFill = fill_elem.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "F0F0F0"})
    solidFill.append(srgbClr)

    ln = fill_elem.makeelement(qn("a:ln"), {"w": "12700"})
    solidFillLn = ln.makeelement(qn("a:solidFill"), {})
    srgbClrLn = solidFillLn.makeelement(qn("a:srgbClr"), {"val": "CCCCCC"})
    solidFillLn.append(srgbClrLn)
    dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    ln.append(solidFillLn)
    ln.append(dash)
    fill_elem.append(ln)

    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = f"📷 {label}"
    p.font.size = Pt(16)
    p.font.color.rgb = MID_GRAY
    p.font.name = "微软雅黑"

    p2 = tf.add_paragraph()
    p2.text = "（请在此处粘贴实际截图）"
    p2.font.size = Pt(12)
    p2.font.color.rgb = MID_GRAY
    p2.font.name = "微软雅黑"
    p2.alignment = PP_ALIGN.CENTER


def _add_table(slide, left, top, width, rows_data: list[list[str]],
               col_widths: list[int] | None = None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 2
    row_height = Inches(0.35)
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "微软雅黑"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if row_idx == 0:
                from pptx.oxml.ns import qn
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn("a:solidFill"), {})
                srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": "2B579A"})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    blank_layout = prs.slide_layouts[6]

    # ========== Slide 1: 封面 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, DARK_BLUE)
    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.0),
                 "基于 Qwen3.5-9B-Base 的动环监控", font_size=36, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.8),
                 "告警诊断场景 LoRA-SFT 训练介绍", font_size=32, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.5),
                 "训练框架：ms-swift  |  训练方式：LoRA-SFT  |  数据：动环监控关联告警", font_size=16,
                 color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                 "ai_trains 训练对比实验平台", font_size=14,
                 color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

    # ========== Slide 2: 目录 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "目录")
    toc_items = [
        "一、项目背景与业务场景",
        "二、模型选型：Qwen3.5-9B-Base",
        "三、数据准备：告警诊断 SFT 数据集",
        "四、环境搭建：ms-swift 安装配置",
        "五、训练配置详解：LoRA-SFT 参数",
        "六、训练执行：swift sft 命令",
        "七、模型导出：LoRA 合并与量化",
        "八、模型部署：Ollama 注册与推理",
        "九、效果评测：离线评测与对比报告",
        "十、平台集成：实验台全链路闭环",
        "十一、效果展示与总结",
    ]
    _add_bullet_points(slide, Inches(1.0), Inches(1.2), Inches(10), Inches(5.5),
                       toc_items, font_size=16, color=DARK_GRAY)

    # ========== Slide 3: 项目背景 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "一、项目背景与业务场景")
    _add_subtitle(slide, "动环监控平台 × 大模型告警诊断")
    items = [
        "业务痛点：动环监控平台每日产生海量告警（UPS、空调、市电、温湿度、门磁等），人工分析效率低、易遗漏",
        "解决思路：利用大模型对关联告警进行自动根因分析，输出结构化诊断报告",
        "涉及设备：UPS系统、开关电源、专用空调、冷水空调、市电配电柜、温湿度传感器、门磁、FSU 等",
        "数据来源：真实动环监控平台关联告警数据，由 Qwen3-235B-A22B 标注生成高质量训练样本",
        "目标：训练 9B 参数量级的 Base 模型，使其具备专业级告警诊断能力",
    ]
    _add_bullet_points(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(3.5),
                       items, font_size=13)
    _add_screenshot_placeholder(slide, Inches(0.8), Inches(4.6), Inches(5), Inches(2.0),
                                "动环监控平台告警界面截图")
    _add_screenshot_placeholder(slide, Inches(6.5), Inches(4.6), Inches(5.5), Inches(2.0),
                                "告警诊断分析报告示例截图")

    # ========== Slide 4: 训练数据样例 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "一、业务场景：告警诊断数据样例")
    _add_subtitle(slide, "输入：关联告警摘要  →  输出：根因分析报告")
    _add_code_block(slide, Inches(0.4), Inches(1.5), Inches(5.5), Inches(5.0), [
        ("// 输入（压缩后的告警摘要）", CODE_COMMENT),
        ("告警摘要（已压缩整理）", CODE_WHITE),
        ("告警窗口：2025-04-07 09:11 至 17:11", CODE_WHITE),
        ("原始告警数：680", CODE_YELLOW),
        ("去重后事件数：42", CODE_YELLOW),
        ("涉及设备数：8", CODE_YELLOW),
        ("高频设备：2#UPS电池组 x320；空调电表 x95", CODE_GREEN),
        ("", CODE_WHITE),
        ("首发关键事件：", CODE_GREEN),
        ("1. 市电停电模块 C相电压过高告警", CODE_WHITE),
        ("2. 1#门磁 门磁状态有告警", CODE_WHITE),
        ("高频重复事件：", CODE_GREEN),
        ("1. 2#UPS系统-4#电池组 单体18电压过高", CODE_WHITE),
        ("2. 2#UPS系统-4#电池组 单体03电压过低", CODE_WHITE),
    ])
    _add_code_block(slide, Inches(6.2), Inches(1.5), Inches(6.5), Inches(5.0), [
        ("// 输出（结构化诊断报告）", CODE_COMMENT),
        ("# 关联告警根因分析报告", CODE_GREEN),
        ("", CODE_WHITE),
        ("## 根因定位", CODE_YELLOW),
        ("根本原因：市电供电异常导致门磁系统误告警", CODE_WHITE),
        ("推理链条：", CODE_GREEN),
        ("1. 市电C相电压过高 → UPS切换响应", CODE_WHITE),
        ("2. 电压波动 → 门磁信号线路电磁干扰", CODE_WHITE),
        ("3. 控制器供电波动 → 状态检测异常", CODE_WHITE),
        ("", CODE_WHITE),
        ("## 影响范围", CODE_YELLOW),
        ("- 安防系统可靠性下降", CODE_WHITE),
        ("- 暴露供电系统抗扰动能力不足", CODE_WHITE),
        ("", CODE_WHITE),
        ("## 解决方案", CODE_YELLOW),
        ("- 紧急：检查UPS供电质量，重启门磁控制器", CODE_WHITE),
        ("- 长期：加装精密稳压模块，独立UPS供电", CODE_WHITE),
    ])

    # ========== Slide 5: 模型选型 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "二、模型选型：Qwen3.5-9B-Base")
    _add_subtitle(slide, "为什么选择 Qwen3.5-9B-Base？")
    _add_table(slide, Inches(0.5), Inches(1.6), Inches(12), [
        ["维度", "说明"],
        ["模型规模", "9B 参数量，单卡 24GB 可完成 LoRA 训练，性价比高"],
        ["Base vs Instruct", "Base 模型可塑性更强，SFT 后诊断风格更贴合专业场景"],
        ["中文能力", "Qwen 系列在中文理解和生成上表现优异，适合中文告警文本"],
        ["长文本支持", "支持 32K+ 上下文，可处理大量关联告警输入"],
        ["开源许可", "Apache 2.0，支持商业使用和二次训练分发"],
        ["ms-swift 兼容", "Qwen3 系列在 ms-swift 中有原生 template 支持"],
    ])
    _add_textbox(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
                 "模型 ID：Qwen/Qwen3.5-9B-Base    模板：qwen3    训练框架：ms-swift (LoRA-SFT)",
                 font_size=14, bold=True, color=MID_BLUE)

    # ========== Slide 6: 数据准备 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "三、数据准备：告警诊断 SFT 数据集")
    _add_subtitle(slide, "从原始告警到训练数据的完整流水线")

    _add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.4),
                 "数据加工流程：", font_size=14, bold=True, color=DARK_GRAY)
    _add_textbox(slide, Inches(0.5), Inches(1.9), Inches(12), Inches(0.6),
                 "原始告警JSON → 告警压缩整理 → 根因提取 → SFT样本构建 → 训练/评测集切分",
                 font_size=13, color=MID_BLUE)

    _add_code_block(slide, Inches(0.4), Inches(2.5), Inches(6.0), Inches(2.0), [
        ("# 执行数据准备脚本", CODE_COMMENT),
        ("python scripts/prepare_alarm_analysis_datasets.py \\", CODE_GREEN),
        ("  --source resources/alarm_analysis.json \\", CODE_WHITE),
        ("  --datasets-root runtime/datasets \\", CODE_WHITE),
        ("  --swift-train-version alarm_analysis_swift_sft_v1 \\", CODE_YELLOW),
        ("  --eval-version alarm_analysis_eval_v1 \\", CODE_WHITE),
        ("  --eval-count 64 --seed 42", CODE_WHITE),
    ])

    _add_textbox(slide, Inches(6.8), Inches(2.5), Inches(5.5), Inches(0.3),
                 "ms-swift JSONL 格式（messages 聊天风格）：", font_size=12, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(6.8), Inches(2.8), Inches(5.8), Inches(1.8), [
        ("{", CODE_WHITE),
        ('  "messages": [', CODE_WHITE),
        ('    {"role": "system",', CODE_YELLOW),
        ('     "content": "你是机房设备运维专家..."},', CODE_WHITE),
        ('    {"role": "user",', CODE_GREEN),
        ('     "content": "告警摘要（已压缩整理）..."},', CODE_WHITE),
        ('    {"role": "assistant",', CODE_YELLOW),
        ('     "content": "# 关联告警根因分析报告..."}', CODE_WHITE),
        ("  ]", CODE_WHITE),
        ("}", CODE_WHITE),
    ])

    _add_table(slide, Inches(0.4), Inches(4.8), Inches(12), [
        ["数据集版本", "格式", "样本数", "用途"],
        ["alarm_analysis_swift_sft_v1", "JSONL (messages)", "~250+", "ms-swift LoRA-SFT 训练"],
        ["alarm_analysis_train_v1", "JSON (Alpaca)", "~250+", "LLaMA-Factory 训练"],
        ["alarm_analysis_eval_v1", "JSONL", "64", "Ollama 离线评测"],
    ])

    # ========== Slide 7: 环境搭建 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "四、环境搭建：ms-swift 安装配置")
    _add_subtitle(slide, "WSL / Linux 环境下的 ms-swift 安装")

    _add_code_block(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(3.5), [
        ("# 1. 创建 conda 环境", CODE_COMMENT),
        ("conda create -n swift python=3.11 -y", CODE_GREEN),
        ("conda activate swift", CODE_GREEN),
        ("", CODE_WHITE),
        ("# 2. 安装 PyTorch (CUDA 12.1)", CODE_COMMENT),
        ("pip install torch torchvision torchaudio \\", CODE_GREEN),
        ("    --index-url https://download.pytorch.org/whl/cu121", CODE_WHITE),
        ("", CODE_WHITE),
        ("# 3. 安装 ms-swift", CODE_COMMENT),
        ("pip install 'ms-swift[llm]' -U", CODE_GREEN),
        ("", CODE_WHITE),
        ("# 4. 验证安装", CODE_COMMENT),
        ("swift sft --help", CODE_GREEN),
        ("python -c \"import torch; print(torch.cuda.is_available())\"", CODE_YELLOW),
    ])

    _add_code_block(slide, Inches(6.5), Inches(1.5), Inches(6.0), Inches(3.5), [
        ("# 5. 下载基础模型（使用 ModelScope）", CODE_COMMENT),
        ("export USE_MODELSCOPE_HUB=1", CODE_GREEN),
        ("# 或使用 HuggingFace", CODE_COMMENT),
        ("export HF_ENDPOINT=https://hf-mirror.com", CODE_YELLOW),
        ("", CODE_WHITE),
        ("# 模型将自动在首次训练时下载", CODE_COMMENT),
        ("# 推荐手动下载到本地缓存目录:", CODE_COMMENT),
        ("modelscope download \\", CODE_GREEN),
        ("  --model Qwen/Qwen3.5-9B-Base \\", CODE_WHITE),
        ("  --local_dir /data/models/Qwen3.5-9B-Base", CODE_WHITE),
        ("", CODE_WHITE),
        ("# 6. GPU 要求", CODE_COMMENT),
        ("# 单卡: NVIDIA A100/A800 40GB+ (推荐)", CODE_YELLOW),
        ("# 单卡: RTX 4090 24GB (需优化 batch_size)", CODE_YELLOW),
    ])

    _add_textbox(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.4),
                 "硬件要求：单卡 24GB+ 显存 (LoRA 训练) | 推荐 A100 40GB 或 RTX 4090",
                 font_size=13, bold=True, color=ACCENT_ORANGE)
    _add_textbox(slide, Inches(0.5), Inches(5.7), Inches(12), Inches(0.8),
                 "软件环境：Python 3.11 + PyTorch 2.x + CUDA 12.1+ + ms-swift latest",
                 font_size=13, color=MID_GRAY)

    # ========== Slide 8: 训练配置详解 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "五、训练配置详解：LoRA-SFT 参数")
    _add_subtitle(slide, "swift.train.yaml 配置文件解析")

    _add_code_block(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(4.8), [
        ("# swift.train.yaml", CODE_COMMENT),
        ("# 基础模型配置", CODE_COMMENT),
        ("model: Qwen/Qwen3.5-9B-Base", CODE_GREEN),
        ("template: qwen3", CODE_YELLOW),
        ("", CODE_WHITE),
        ("# 数据集配置", CODE_COMMENT),
        ("dataset:", CODE_GREEN),
        ("  - runtime/datasets/alarm_analysis_swift_sft_v1/", CODE_WHITE),
        ("    alarm_analysis_swift_sft_v1.jsonl", CODE_WHITE),
        ("", CODE_WHITE),
        ("# LoRA 参数", CODE_COMMENT),
        ("train_type: lora", CODE_GREEN),
        ("lora_rank: 64", CODE_YELLOW),
        ("lora_alpha: 128", CODE_YELLOW),
        ("target_modules: all-linear", CODE_GREEN),
        ("", CODE_WHITE),
        ("# 训练超参", CODE_COMMENT),
        ("num_train_epochs: 3", CODE_YELLOW),
        ("max_length: 4096", CODE_WHITE),
        ("per_device_train_batch_size: 1", CODE_WHITE),
        ("gradient_accumulation_steps: 8", CODE_WHITE),
        ("learning_rate: 2.0e-4", CODE_YELLOW),
    ])

    _add_table(slide, Inches(6.5), Inches(1.5), Inches(6.0), [
        ["参数", "值", "说明"],
        ["model", "Qwen/Qwen3.5-9B-Base", "基础模型路径"],
        ["train_type", "lora", "LoRA 低秩适应"],
        ["lora_rank", "64", "LoRA 秩（能力与显存平衡）"],
        ["lora_alpha", "128", "LoRA 缩放系数（通常=2×rank）"],
        ["target_modules", "all-linear", "对所有线性层应用 LoRA"],
        ["template", "qwen3", "Qwen3 对话模板"],
        ["max_length", "4096", "最大序列长度"],
        ["num_train_epochs", "3", "训练轮数"],
        ["batch_size", "1", "单卡 batch size"],
        ["grad_accum", "8", "梯度累积步数（等效 bs=8）"],
        ["learning_rate", "2e-4", "学习率"],
    ])

    # ========== Slide 9: 训练命令 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "六、训练执行：swift sft 命令")
    _add_subtitle(slide, "一键启动 LoRA-SFT 训练")

    _add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.3),
                 "方式一：使用配置文件（推荐）", font_size=14, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(0.4), Inches(1.8), Inches(12.2), Inches(1.2), [
        ("# 使用 YAML 配置文件启动训练", CODE_COMMENT),
        ("CUDA_VISIBLE_DEVICES=0 swift sft --config swift.train.yaml", CODE_GREEN),
    ])

    _add_textbox(slide, Inches(0.5), Inches(3.2), Inches(12), Inches(0.3),
                 "方式二：命令行直接指定参数", font_size=14, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(0.4), Inches(3.5), Inches(12.2), Inches(2.5), [
        ("CUDA_VISIBLE_DEVICES=0 swift sft \\", CODE_GREEN),
        ("  --model Qwen/Qwen3.5-9B-Base \\", CODE_WHITE),
        ("  --dataset runtime/datasets/alarm_analysis_swift_sft_v1/"
         "alarm_analysis_swift_sft_v1.jsonl \\", CODE_WHITE),
        ("  --train_type lora \\", CODE_YELLOW),
        ("  --template qwen3 \\", CODE_WHITE),
        ("  --max_length 4096 \\", CODE_WHITE),
        ("  --num_train_epochs 3 \\", CODE_WHITE),
        ("  --per_device_train_batch_size 1 \\", CODE_WHITE),
        ("  --gradient_accumulation_steps 8 \\", CODE_WHITE),
        ("  --lora_rank 64 \\", CODE_YELLOW),
        ("  --lora_alpha 128 \\", CODE_YELLOW),
        ("  --target_modules all-linear \\", CODE_WHITE),
        ("  --output_dir runtime/runs/{run_id}/output/adapter", CODE_GREEN),
    ])

    _add_screenshot_placeholder(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
                                "训练过程 loss 曲线截图（训练完成后粘贴）")

    # ========== Slide 10: 多卡训练 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "六、进阶：多卡训练与 DeepSpeed")
    _add_subtitle(slide, "加速训练的多种方式")

    _add_textbox(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.3),
                 "多卡数据并行训练：", font_size=14, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(0.4), Inches(1.8), Inches(6.0), Inches(1.4), [
        ("# 双卡数据并行", CODE_COMMENT),
        ("CUDA_VISIBLE_DEVICES=0,1 \\", CODE_GREEN),
        ("NPROC_PER_NODE=2 \\", CODE_GREEN),
        ("swift sft --config swift.train.yaml", CODE_YELLOW),
    ])

    _add_textbox(slide, Inches(0.5), Inches(3.5), Inches(6), Inches(0.3),
                 "使用 DeepSpeed ZeRO-2 优化显存：", font_size=14, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(0.4), Inches(3.8), Inches(6.0), Inches(1.6), [
        ("# DeepSpeed ZeRO-2（推荐多卡场景）", CODE_COMMENT),
        ("CUDA_VISIBLE_DEVICES=0,1 \\", CODE_GREEN),
        ("NPROC_PER_NODE=2 \\", CODE_GREEN),
        ("swift sft \\", CODE_YELLOW),
        ("  --config swift.train.yaml \\", CODE_WHITE),
        ("  --deepspeed zero2", CODE_GREEN),
    ])

    _add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.3),
                 "训练监控（TensorBoard/SwanLab）：", font_size=14, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(6.7), Inches(1.8), Inches(5.8), Inches(1.6), [
        ("# 查看训练日志", CODE_COMMENT),
        ("tensorboard --logdir runtime/runs/{run_id}/logs", CODE_GREEN),
        ("", CODE_WHITE),
        ("# 或使用 SwanLab", CODE_COMMENT),
        ("swift sft --config swift.train.yaml \\", CODE_YELLOW),
        ("  --report_to swanlab", CODE_GREEN),
    ])

    _add_screenshot_placeholder(slide, Inches(6.7), Inches(3.8), Inches(5.8), Inches(2.5),
                                "TensorBoard 训练 loss 曲线截图")

    # ========== Slide 11: 模型导出 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "七、模型导出：LoRA 合并与量化")
    _add_subtitle(slide, "训练完成后的模型导出流程")

    _add_textbox(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(0.3),
                 "步骤一：合并 LoRA 适配器到基础模型", font_size=14, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(0.4), Inches(1.8), Inches(12.2), Inches(1.5), [
        ("# 合并 LoRA adapter 到基础模型", CODE_COMMENT),
        ("swift export \\", CODE_GREEN),
        ("  --model Qwen/Qwen3.5-9B-Base \\", CODE_WHITE),
        ("  --adapters runtime/runs/{run_id}/output/adapter \\", CODE_YELLOW),
        ("  --merge_lora true \\", CODE_WHITE),
        ("  --output_dir runtime/runs/{run_id}/output/merged", CODE_GREEN),
    ])

    _add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12), Inches(0.3),
                 "步骤二：量化为 GGUF 格式（用于 Ollama 部署）", font_size=14, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(0.4), Inches(3.9), Inches(12.2), Inches(1.5), [
        ("# 使用 llama.cpp 转换为 GGUF 格式", CODE_COMMENT),
        ("python llama.cpp/convert_hf_to_gguf.py \\", CODE_GREEN),
        ("  runtime/runs/{run_id}/output/merged \\", CODE_WHITE),
        ("  --outfile runtime/artifacts/gguf/alarm-qwen35-9b-sft.gguf \\", CODE_YELLOW),
        ("  --outtype q4_k_m", CODE_GREEN),
    ])

    _add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12), Inches(0.5),
                 "导出产物：merged/ (完整 HF 权重) → .gguf (量化后) → Ollama 模型",
                 font_size=13, bold=True, color=ACCENT_ORANGE)

    _add_screenshot_placeholder(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
                                "导出命令执行输出截图")

    # ========== Slide 12: Ollama 部署 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "八、模型部署：Ollama 注册与推理")
    _add_subtitle(slide, "本地推理服务搭建")

    _add_code_block(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(2.5), [
        ("# 1. 编写 Modelfile", CODE_COMMENT),
        ('FROM ./alarm-qwen35-9b-sft.gguf', CODE_GREEN),
        ("", CODE_WHITE),
        ('SYSTEM """你是机房设备运维专家，请基于', CODE_YELLOW),
        ('压缩整理后的动环监控平台关联告警摘要，', CODE_YELLOW),
        ('分析根本原因、影响范围和处置建议。"""', CODE_YELLOW),
        ("", CODE_WHITE),
        ("PARAMETER temperature 0.1", CODE_WHITE),
        ("PARAMETER top_p 0.8", CODE_WHITE),
        ("PARAMETER num_ctx 8192", CODE_GREEN),
        ("PARAMETER seed 42", CODE_WHITE),
    ])

    _add_code_block(slide, Inches(6.5), Inches(1.5), Inches(6.0), Inches(2.5), [
        ("# 2. 注册到 Ollama", CODE_COMMENT),
        ("ollama create alarm-swift-qwen35-9b \\", CODE_GREEN),
        ("  -f Modelfile", CODE_WHITE),
        ("", CODE_WHITE),
        ("# 3. 测试推理", CODE_COMMENT),
        ("ollama run alarm-swift-qwen35-9b \\", CODE_GREEN),
        ('  "告警摘要（已压缩整理）\\n', CODE_WHITE),
        ('   告警窗口：2025-04-07 09:11至17:11\\n', CODE_WHITE),
        ('   原始告警数：680\\n', CODE_WHITE),
        ('   涉及设备：UPS,空调,市电..."', CODE_WHITE),
    ])

    _add_textbox(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(0.3),
                 "平台自动化：Worker 自动完成 export → GGUF → Ollama 注册 → 评测", font_size=13,
                 bold=True, color=MID_BLUE)

    _add_screenshot_placeholder(slide, Inches(0.5), Inches(4.7), Inches(5.5), Inches(2.0),
                                "Ollama 推理输出截图")
    _add_screenshot_placeholder(slide, Inches(6.5), Inches(4.7), Inches(5.5), Inches(2.0),
                                "平台 H5 页面推理效果截图")

    # ========== Slide 13: 评测 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "九、效果评测：离线评测与对比报告")
    _add_subtitle(slide, "基于 Ollama 的统一评测框架")

    _add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.0), [
        "评测集：alarm_analysis_eval_v1（64 条）",
        "Benchmark 集：16 条短文本（压测吞吐）",
        "Smoke 集：8 条（快速冒烟验证）",
        "评测指标：expected_contains 关键词命中率",
        "自动化评测：Worker 一键完成",
    ], font_size=13)

    _add_code_block(slide, Inches(6.3), Inches(1.5), Inches(6.2), Inches(2.2), [
        ("# 评测数据格式 (eval.jsonl)", CODE_COMMENT),
        ("{", CODE_WHITE),
        ('  "task_type": "alarm_diagnosis",', CODE_YELLOW),
        ('  "messages": [...],', CODE_WHITE),
        ('  "expected_contains": [', CODE_GREEN),
        ('    "市电供电异常导致门磁系统误告警",', CODE_WHITE),
        ('    "UPS", "市电"', CODE_WHITE),
        ("  ]", CODE_WHITE),
        ("}", CODE_WHITE),
    ])

    _add_textbox(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.3),
                 "多模型对比矩阵：", font_size=14, bold=True, color=DARK_GRAY)
    _add_table(slide, Inches(0.4), Inches(4.2), Inches(12), [
        ["实验名称", "基础模型", "训练方式", "评测命中率", "推理延迟"],
        ["基线 (原始推理)", "Qwen3.5-9B-Base", "无训练", "待填写", "待填写"],
        ["LoRA-SFT (swift)", "Qwen3.5-9B-Base", "LoRA + SFT", "待填写", "待填写"],
        ["LoRA-SFT (LLaMA-Factory)", "Qwen3.5-2B", "LoRA + SFT", "待填写", "待填写"],
    ])

    _add_screenshot_placeholder(slide, Inches(0.5), Inches(5.8), Inches(5.5), Inches(0.9),
                                "评测结果对比图表截图")
    _add_screenshot_placeholder(slide, Inches(6.5), Inches(5.8), Inches(5.5), Inches(0.9),
                                "对比报告 H5 页面截图")

    # ========== Slide 14: 平台集成 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "十、平台集成：实验台全链路闭环")
    _add_subtitle(slide, "ai_trains 训练对比实验平台架构")

    _add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(3.5), [
        "前端：Vue3 + Vite H5 实验管理界面",
        "后端：FastAPI + SQLite 元数据管理",
        "Worker：轮询任务，串行执行训练流程",
        "训练后端：LLaMA-Factory (一期) + ms-swift (二期)",
        "推理后端：Ollama 统一离线评测",
        "报告系统：单次运行报告 + 多模型对比报告",
        "",
        "全链路流程：",
        "  创建实验 → 生成运行 → 准备配置",
        "  → 执行训练 → 导出模型 → 注册Ollama",
        "  → 离线评测 → 生成报告",
    ], font_size=12)

    _add_textbox(slide, Inches(6.5), Inches(1.5), Inches(5.5), Inches(0.3),
                 "创建 swift 训练实验 (API 请求)：", font_size=12, bold=True, color=DARK_GRAY)
    _add_code_block(slide, Inches(6.4), Inches(1.8), Inches(6.0), Inches(3.2), [
        ("POST /api/experiments", CODE_COMMENT),
        ("{", CODE_WHITE),
        ('  "name": "Qwen3.5-9B-Base 告警诊断",', CODE_YELLOW),
        ('  "base_model": "Qwen/Qwen3.5-9B-Base",', CODE_GREEN),
        ('  "trainer_backend": "swift",', CODE_GREEN),
        ('  "route_type": "sft",', CODE_WHITE),
        ('  "dataset_version":', CODE_WHITE),
        ('    "alarm_analysis_swift_sft_v1",', CODE_YELLOW),
        ('  "train_config": {', CODE_WHITE),
        ('    "train_type": "lora",', CODE_WHITE),
        ('    "template": "qwen3",', CODE_WHITE),
        ('    "lora_rank": 64,', CODE_YELLOW),
        ('    "lora_alpha": 128', CODE_YELLOW),
        ("  }", CODE_WHITE),
        ("}", CODE_WHITE),
    ])

    _add_screenshot_placeholder(slide, Inches(0.5), Inches(5.3), Inches(5.5), Inches(1.4),
                                "实验台 H5 主界面截图")
    _add_screenshot_placeholder(slide, Inches(6.5), Inches(5.3), Inches(5.5), Inches(1.4),
                                "运行详情页面截图")

    # ========== Slide 15: 训练产物 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "十、运行产物与目录结构")
    _add_subtitle(slide, "每次训练运行生成的文件清单")

    _add_code_block(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(4.5), [
        ("runtime/runs/run_20260412_001/", CODE_GREEN),
        ("├── input/", CODE_WHITE),
        ("│   ├── run_manifest.json", CODE_WHITE),
        ("│   ├── swift.train.yaml      # 训练配置", CODE_YELLOW),
        ("│   ├── swift.export.yaml     # 导出配置", CODE_YELLOW),
        ("│   ├── swift.eval.yaml       # 评测计划", CODE_YELLOW),
        ("│   └── swift.preview.txt", CODE_WHITE),
        ("├── output/", CODE_WHITE),
        ("│   ├── adapter/              # LoRA 权重", CODE_GREEN),
        ("│   ├── merged/               # 合并后模型", CODE_GREEN),
        ("│   ├── Modelfile             # Ollama 配置", CODE_WHITE),
        ("│   └── command_preview.json", CODE_WHITE),
        ("├── logs/", CODE_WHITE),
        ("│   ├── train.log", CODE_WHITE),
        ("│   ├── export.log", CODE_WHITE),
        ("│   └── worker.log", CODE_WHITE),
        ("├── eval/", CODE_WHITE),
        ("│   ├── eval_results.jsonl", CODE_YELLOW),
        ("│   └── eval_summary.json", CODE_YELLOW),
        ("├── benchmark/", CODE_WHITE),
        ("└── reports/", CODE_WHITE),
    ])

    _add_table(slide, Inches(6.5), Inches(1.5), Inches(6.0), [
        ["产物", "路径", "说明"],
        ["训练配置", "input/swift.train.yaml", "完整训练参数"],
        ["LoRA 权重", "output/adapter/", "训练产出的适配器"],
        ["合并模型", "output/merged/", "合并后完整模型"],
        ["训练日志", "logs/train.log", "训练过程详细日志"],
        ["评测结果", "eval/eval_results.jsonl", "逐条评测输出"],
        ["评测汇总", "eval/eval_summary.json", "命中率等指标"],
        ["运行报告", "reports/{run_id}.json", "完整运行报告"],
    ])

    # ========== Slide 16: 效果展示 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, WHITE)
    _add_title(slide, "十一、效果展示")
    _add_subtitle(slide, "训练前后对比")

    _add_screenshot_placeholder(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(2.2),
                                "训练前：Base 模型原始推理输出截图")
    _add_screenshot_placeholder(slide, Inches(6.5), Inches(1.5), Inches(5.8), Inches(2.2),
                                "训练后：SFT 模型诊断输出截图")

    _add_screenshot_placeholder(slide, Inches(0.4), Inches(4.0), Inches(5.8), Inches(2.5),
                                "训练 Loss 曲线截图")
    _add_screenshot_placeholder(slide, Inches(6.5), Inches(4.0), Inches(5.8), Inches(2.5),
                                "多模型对比报告截图")

    # ========== Slide 17: 总结 ==========
    slide = prs.slides.add_slide(blank_layout)
    _add_bg(slide, DARK_BLUE)
    _add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
                 "总结与展望", font_size=32, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    summary_items = [
        "采用 Qwen3.5-9B-Base 作为基础模型，面向动环监控告警诊断场景进行 LoRA-SFT 训练",
        "使用 ms-swift 框架，支持 YAML 配置与命令行两种训练方式",
        "数据集来源于真实动环平台，经 Qwen3-235B 标注 + 自动压缩清洗",
        "全链路自动化：训练 → 导出 → Ollama 部署 → 离线评测 → 对比报告",
        "平台层双后端设计：LLaMA-Factory (一期) + ms-swift (二期) 无缝切换",
    ]
    _add_bullet_points(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5),
                       summary_items, font_size=15, color=WHITE)

    _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
                 "后续计划：", font_size=18, bold=True, color=RGBColor(0xFF, 0xCC, 0x00))
    next_items = [
        "WSL 环境部署 ms-swift 完成真实训练闭环",
        "引入 DAPT/CPT（领域自适应预训练）提升专业术语理解",
        "扩展评测维度：诊断准确率、根因命中率、推理时效",
        "接入更多设备类型和告警场景数据",
    ]
    _add_bullet_points(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0),
                       next_items, font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF))

    prs.save(str(OUTPUT_PATH))
    print(f"PPT saved to: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
